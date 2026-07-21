"""
pptx 内置技能 — PowerPoint 演示文稿 (.pptx) 的创建、编辑和读取。
支持幻灯片、文本框、图片、表格和基础格式化。
"""
from typing import Any, Optional
from loguru import logger

SKILL_NAME = "pptx"
SKILL_DESCRIPTION = "创建、编辑和读取 PowerPoint 演示文稿 (.pptx)，支持幻灯片/文本/图片/表格操作"

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


async def execute(
    action: str,
    file_path: Optional[str] = None,
    content: Optional[str] = None,
    slide_index: Optional[int] = None,
    title: Optional[str] = None,
    **kwargs: Any,
) -> dict[str, Any]:
    """
    执行 PowerPoint 文档操作。

    Args:
        action: 操作类型（read/create/add_slide/add_text/add_image/info）
        file_path: 文档路径
        content: 文本内容
        slide_index: 幻灯片索引（0-based）
        title: 幻灯片标题

    Returns:
        操作结果
    """
    if not HAS_PPTX:
        return {
            "success": False,
            "error": "缺少 python-pptx 依赖，请运行: pip install python-pptx",
        }

    valid_actions = {"read", "create", "add_slide", "add_text", "add_image", "info"}
    if action not in valid_actions:
        return {"success": False, "error": f"不支持的操作: {action}，支持: {', '.join(sorted(valid_actions))}"}

    # ---- 读取演示文稿 ----
    if action == "read":
        if not file_path:
            return {"success": False, "error": "读取操作需要提供 file_path"}
        try:
            prs = Presentation(file_path)
            slides_info = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                slides_info.append({
                    "index": i,
                    "shape_count": len(slide.shapes),
                    "text_preview": " ".join(texts)[:200] if texts else "",
                })
            return {
                "success": True,
                "action": "read",
                "file": file_path,
                "slide_count": len(prs.slides),
                "slides": slides_info,
            }
        except Exception as e:
            return {"success": False, "error": f"读取失败: {str(e)}"}

    # ---- 创建新演示文稿 ----
    elif action == "create":
        if not file_path:
            return {"success": False, "error": "创建需要提供 file_path"}
        try:
            prs = Presentation()
            # 添加标题幻灯片
            title_text = title or "演示文稿"
            subtitle_text = content or ""
            slide_layout = prs.slide_layouts[0]  # 标题幻灯片版式
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle_text
            prs.save(file_path)
            return {
                "success": True,
                "action": "create",
                "file": file_path,
                "slide_count": 1,
            }
        except Exception as e:
            return {"success": False, "error": f"创建失败: {str(e)}"}

    # ---- 添加幻灯片 ----
    elif action == "add_slide":
        if not file_path:
            return {"success": False, "error": "add_slide 需要提供 file_path"}
        try:
            prs = Presentation(file_path)
            slide_layout = prs.slide_layouts[1]  # 标题和内容版式
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title or ""
            if content and len(slide.placeholders) > 1:
                slide.placeholders[1].text = content
            prs.save(file_path)
            return {
                "success": True,
                "action": "add_slide",
                "file": file_path,
                "slide_index": len(prs.slides) - 1,
                "slide_count": len(prs.slides),
            }
        except Exception as e:
            return {"success": False, "error": f"添加幻灯片失败: {str(e)}"}

    # ---- 向幻灯片添加文本 ----
    elif action == "add_text":
        if not file_path or content is None:
            return {"success": False, "error": "add_text 需要提供 file_path 和 content"}
        try:
            prs = Presentation(file_path)
            idx = slide_index or (len(prs.slides) - 1)
            if idx < 0 or idx >= len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {idx}，共 {len(prs.slides)} 张"}
            slide = prs.slides[idx]
            # 添加文本框
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            for line in content.split("\n"):
                p = tf.add_paragraph()
                p.text = line
            prs.save(file_path)
            return {
                "success": True,
                "action": "add_text",
                "file": file_path,
                "slide_index": idx,
            }
        except Exception as e:
            return {"success": False, "error": f"添加文本失败: {str(e)}"}

    # ---- 文档信息 ----
    elif action == "info":
        if not file_path:
            return {"success": False, "error": "info 操作需要提供 file_path"}
        try:
            prs = Presentation(file_path)
            return {
                "success": True,
                "file": file_path,
                "slide_count": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
            }
        except Exception as e:
            return {"success": False, "error": f"读取信息失败: {str(e)}"}

    return {"success": False, "error": "未识别的操作"}
